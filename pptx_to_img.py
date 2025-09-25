# Copyright (c) OpenAI. All rights reserved.
import argparse
import os
import subprocess
import tempfile
from typing import Any, Sequence, cast

import numpy as np
from pdf2image import convert_from_path
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.util import Emu

# Overflow checker configuration
PAD_PX: int = 100  # fixed padding on every side in pixels
EMU_PER_INCH: int = 914_400
PAD_RGB = (200, 200, 200)


def calc_dpi(prs: Any, max_w_px: int, max_h_px: int) -> int:
    """Calculate DPI so that the rendered slide fits within the given box."""
    width_in = prs.slide_width / EMU_PER_INCH
    height_in = prs.slide_height / EMU_PER_INCH
    return round(min(max_w_px / width_in, max_h_px / height_in))


def rasterize(pptx_path: str, out_dir: str, dpi: int) -> Sequence[str]:
    """Rasterise PPTX to PNG files placed in *out_dir* and return the image paths."""

    os.makedirs(out_dir, exist_ok=True)

    pptx_path = os.path.abspath(pptx_path)
    work_dir = os.path.dirname(pptx_path)

    subprocess.run(
        [
            "soffice",
            "--headless",
            "--convert-to",
            "pdf",
            "--outdir",
            work_dir,
            pptx_path,
        ],
        check=True,
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
    )

    pdf_path = os.path.join(work_dir, f"{os.path.splitext(os.path.basename(pptx_path))[0]}.pdf")

    if not os.path.exists(pdf_path):
        raise RuntimeError("Failed to produce PDF for overflow detection.")

    paths_raw = cast(
        list[str],
        convert_from_path(
            pdf_path,
            dpi=dpi,
            fmt="png",
            thread_count=8,
            output_folder=out_dir,
            paths_only=True,
            output_file="slide",
        ),
    )
    # Rename convert_from_path's output format f'slide{thread_id:04d}-{page_num:02d}.png'
    slides = []
    for src_path in paths_raw:
        base = os.path.splitext(os.path.basename(src_path))[0]
        slide_num_str = base.split("-")[-1]
        slide_num = int(slide_num_str)
        dst_path = os.path.join(out_dir, f"slide-{slide_num}.png")
        os.replace(src_path, dst_path)
        slides.append((slide_num, dst_path))
    slides.sort(key=lambda t: t[0])
    final_paths = [path for _, path in slides]
    return final_paths


def px_to_emu(px: int, dpi: int) -> Emu:
    return Emu(int(px * EMU_PER_INCH // dpi))


def calc_tol(dpi: int) -> int:
    """Calculate per-channel colour tolerance appropriate for *dpi* (anti-aliasing tolerance)."""
    if dpi >= 300:
        return 0
    # 1 at 250 DPI, 5 at 150 DPI, capped to 10.
    tol = round((300 - dpi) / 25)
    return min(max(tol, 1), 10)


def enlarge_deck(src: str, dst: str, pad_emu: Emu) -> tuple[int, int]:
    """Enlarge the input PPTX with a fixed grey padding and returns the new page size."""
    prs = Presentation(src)
    w0 = cast(Emu, prs.slide_width)
    h0 = cast(Emu, prs.slide_height)
    w1 = Emu(w0 + 2 * pad_emu)
    h1 = Emu(h0 + 2 * pad_emu)
    prs.slide_width = w1
    prs.slide_height = h1

    for slide in prs.slides:
        # Shift all shapes so the original canvas sits centred in the new deck.
        for shp in list(slide.shapes):
            shp.left = Emu(int(shp.left) + pad_emu)
            shp.top = Emu(int(shp.top) + pad_emu)

        pads = (
            (Emu(0), Emu(0), pad_emu, h1),  # left
            (Emu(int(w1) - int(pad_emu)), Emu(0), pad_emu, h1),  # right
            (Emu(0), Emu(0), w1, pad_emu),  # top
            (Emu(0), Emu(int(h1) - int(pad_emu)), w1, pad_emu),  # bottom
        )

        sp_tree = slide.shapes._spTree  # pylint: disable=protected-access

        for left, top, width, height in pads:
            pad_shape = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, width, height
            )
            pad_shape.fill.solid()
            pad_shape.fill.fore_color.rgb = RGBColor(*PAD_RGB)
            pad_shape.line.fill.background()

            # Send pad behind all other shapes (index 2 after mandatory nodes)
            sp_tree.remove(pad_shape._element)
            sp_tree.insert(2, pad_shape._element)

    prs.save(dst)
    return int(w1), int(h1)


def inspect_images(
    paths: Sequence[str],
    pad_ratio_w: float,
    pad_ratio_h: float,
    dpi: int,
) -> list[int]:
    """Return 1-based indices of slides that contain pixels outside the pad."""

    tol = calc_tol(dpi)
    failures: list[int] = []
    pad_colour = np.array(PAD_RGB, dtype=np.uint8)

    for idx, img_path in enumerate(paths, start=1):
        with Image.open(img_path) as img:
            rgb = img.convert("RGB")
            arr = np.asarray(rgb)

        h, w, _ = arr.shape
        # Exclude the innermost 1-pixel band
        pad_x = int(w * pad_ratio_w) - 1
        pad_y = int(h * pad_ratio_h) - 1

        left_margin = arr[:, :pad_x, :]
        right_margin = arr[:, w - pad_x :, :]
        top_margin = arr[:pad_y, :, :]
        bottom_margin = arr[h - pad_y :, :, :]

        def _is_clean(margin: np.ndarray) -> bool:
            diff = np.abs(margin.astype(np.int16) - pad_colour)
            matches = np.all(diff <= tol, axis=-1)
            mismatch_fraction = 1.0 - (np.count_nonzero(matches) / matches.size)
            if dpi >= 300:
                max_mismatch = 0.01
            elif dpi >= 200:
                max_mismatch = 0.02
            else:
                max_mismatch = 0.03
            return mismatch_fraction <= max_mismatch

        if not (
            _is_clean(left_margin)
            and _is_clean(right_margin)
            and _is_clean(top_margin)
            and _is_clean(bottom_margin)
        ):
            failures.append(idx)

    return failures


def check_overflow(pptx_path: str, dpi: int) -> None:
    """Emit a warning if input PPTX contains any edge-overflowing content."""

    # Not using ``tempfile.TemporaryDirectory(delete=False)`` for Python 3.11 compatibility.
    tmpdir = tempfile.mkdtemp()
    enlarged_pptx = os.path.join(tmpdir, "enlarged.pptx")
    pad_emu = px_to_emu(PAD_PX, dpi)
    w1, h1 = enlarge_deck(pptx_path, enlarged_pptx, pad_emu=pad_emu)
    pad_ratio_w = pad_emu / w1
    pad_ratio_h = pad_emu / h1

    img_dir = os.path.join(tmpdir, "imgs")
    img_paths = rasterize(enlarged_pptx, img_dir, dpi)
    failing = inspect_images(img_paths, pad_ratio_w, pad_ratio_h, dpi)

    if failing:
        print(
            "WARNING: Slides with content overflowing original canvas (1-based indexing): "
            + ", ".join(map(str, failing))
            + "\n"
            + "    Rendered images with grey paddings for problematic slides are available at: "
        )
        # Provide full filesystem paths to the rendered images for each failing slide
        for i in failing:
            print("   ", img_paths[i - 1])
        print(
            "    Please also check other slides for potential issues and fix them if there are any."
        )


def main() -> None:
    parser = argparse.ArgumentParser(description="Convert PPTX file to images.")
    parser.add_argument(
        "--input",
        type=str,
        required=True,
        help="Path to the input PPTX file.",
    )
    parser.add_argument(
        "--output",
        type=str,
        default="ppt-preview",
        help="Output directory for the rendered PNGs (default: ppt-preview)",
    )
    parser.add_argument(
        "--width",
        type=int,
        default=1600,
        help="Approximate maximum width in pixels after isotropic scaling (default 1600). The actual value may exceed slightly.",
    )
    parser.add_argument(
        "--height",
        type=int,
        default=900,
        help="Approximate maximum height in pixels after isotropic scaling (default 900). The actual value may exceed slightly.",
    )
    args = parser.parse_args()

    out_dir = os.path.abspath(args.output)

    # TODO: remove after the container's file permission issue is fixed
    # Ensure the output directory has the correct permissions **before** rendering.
    # os.makedirs(out_dir, exist_ok=True)
    # os.chmod(out_dir, 0o770)

    pres = Presentation(args.input)
    dpi = calc_dpi(pres, args.width, args.height)
    check_overflow(args.input, dpi)
    rasterize(args.input, out_dir, dpi)
    print("Saved rendered slides (slide-1.png, slide-2.png, etc.) to " + out_dir)


if __name__ == "__main__":
    main()
